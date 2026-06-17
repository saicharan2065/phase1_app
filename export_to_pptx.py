import pptx
from pptx.util import Inches, Pt

def create_presentation():
    prs = pptx.Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Antigravity OS - Phase 1"
    subtitle.text = "Financial Crime Intelligence & AMD MI300X Multi-Modal Engine Architecture"
    
    # Slide 2: Hardware Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "Hardware Architecture & VRAM Splitting"
    tf = body.text_frame
    tf.text = "Successfully engineered a multi-modal execution pipeline leveraging massive hardware."
    p = tf.add_paragraph()
    p.text = "• Eradicated PyTorch CUDA deadlocks via strict hardware isolation."
    p = tf.add_paragraph()
    p.text = "• Spills massive 70B parameter LLMs into 2TB of System RAM for parallel QLoRA training."
    p = tf.add_paragraph()
    p.text = "• Preserves 192GB VRAM strictly for Bulk SAR Inference and Llava-13B Vision Lab generation."
    
    # Slide 3: The Four Core Engines
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "The 4 Simultaneous MI300X Engines"
    tf = body.text_frame
    tf.text = "All four engines successfully synchronize behind a threading barrier before firing simultaneously:"
    p = tf.add_paragraph()
    p.text = "1. Bulk SAR Generator: Mass-producing Suspicious Activity Reports (Qwen-70B)."
    p = tf.add_paragraph()
    p.text = "2. QLoRA Studio: On-the-fly finetuning isolated on System RAM (SFTTrainer/PEFT)."
    p = tf.add_paragraph()
    p.text = "3. Vision Forensics Lab: Deepfake detection via multi-modal image tokenization (Llava-13B)."
    p = tf.add_paragraph()
    p.text = "4. GNN Topography: Processing 50M+ node graphs for financial risk clustering."
    
    # Slide 4: Real Data & Telemetry
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "Real Data Ingestion & Live Telemetry"
    tf = body.text_frame
    tf.text = "Zero simulations. The platform ingests direct streams from Hugging Face data repositories."
    p = tf.add_paragraph()
    p.text = "• Live Engine Telemetry: Real-time UI polling tracks accurate batch processing down to the microsecond."
    p = tf.add_paragraph()
    p.text = "• Completely dynamic data rendering for text and images natively inside the GPU context."
    p = tf.add_paragraph()
    p.text = "• Bulletproof memory management capable of surviving silent PEFT crashes."
    
    prs.save('Phase1_Pitch_Deck.pptx')
    print("Successfully generated 4-page presentation: Phase1_Pitch_Deck.pptx")

if __name__ == "__main__":
    create_presentation()
